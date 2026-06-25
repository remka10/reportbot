#!/usr/bin/env python3
"""
Скрипт для подготовки PPTX-шаблона отчёта.

Запускать ОДИН РАЗ после первой настройки:
    python scripts/prepare_pptx_template.py

Что делает:
  1. Читает существующий Report_2026-2.pptx (или указанный файл)
  2. Расставляет плейсхолдеры {{student_name}}, {{department_color}} и т.д.
     в нужных текстовых блоках на каждом слайде
  3. Сохраняет результат как app/templates/report_template.pptx

Плейсхолдеры, которые понимает PptxService:
  {{student_name}}      — ФИО ребёнка
  {{shift_name}}        — Название смены
  {{department_name}}   — Название департамента
  {{teacher_name}}      — ФИО педагога
  {{report_date}}       — Дата генерации (дд.мм.гггг)
  {{report_text}}       — Полный текст отчёта (один блок)
  {{department_color}}  — Текстовый фрейм с этим тегом перекрашивается
                          в цвет департамента (сам тег удаляется)
  {{block_title}}       — На слайде-шаблоне блока: название блока
  {{block_content}}     — На слайде-шаблоне блока: текст блока

Слайд с {{block_title}} / {{block_content}} будет дублироваться
автоматически под каждый блок отчёта.

Используй этот скрипт как СПРАВОЧНИК —
вручную открой PPTX в PowerPoint/Keynote и вставь плейсхолдеры
в нужные текстовые поля согласно схеме ниже.
"""

import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Установи зависимость: pip install python-pptx")
    sys.exit(1)


# ─── Настройки ────────────────────────────────────────────────────────────────
SCRIPT_DIR   = Path(__file__).parent
PROJECT_ROOT = SCRIPT_DIR.parent
SOURCE_PPTX  = PROJECT_ROOT / "Report_2026-2.pptx"       # ваш оригинальный шаблон
OUTPUT_PPTX  = PROJECT_ROOT / "app" / "templates" / "report_template.pptx"

# Карта: (индекс слайда 0-based, индекс shape) → плейсхолдер
# Заполни согласно реальной структуре твоего шаблона.
# Запусти скрипт с флагом --inspect для просмотра всех shape-ов.
PLACEHOLDER_MAP: dict[tuple[int, int], str] = {
    # Слайд 0 (обложка)
    (0, 2): "{{student_name}}",
    (0, 3): "{{shift_name}}",
    (0, 4): "{{department_name}}{{department_color}}",  # перекрашивается
    (0, 5): "{{teacher_name}}",
    # Слайд 1 (основной текст)
    (1, 1): "{{report_text}}",
    (1, 2): "{{report_date}}",
}


def inspect_pptx(path: Path) -> None:
    """Выводит структуру всех слайдов и shape-ов — используй для маппинга."""
    prs = Presentation(str(path))
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n{'='*60}")
        print(f"Слайд {slide_idx} ({len(slide.shapes)} shape-ов)")
        print(f"{'='*60}")
        for shape_idx, shape in enumerate(slide.shapes):
            text_preview = ""
            if shape.has_text_frame:
                text_preview = shape.text_frame.text[:80].replace("\n", " ↵ ")
            print(f"  [{shape_idx}] {shape.shape_type!s:20s} name={shape.name!r:30s} text={text_preview!r}")


def apply_placeholders(source: Path, output: Path) -> None:
    """Вставляет плейсхолдеры в шаблон согласно PLACEHOLDER_MAP."""
    prs = Presentation(str(source))

    for (slide_idx, shape_idx), placeholder in PLACEHOLDER_MAP.items():
        if slide_idx >= len(prs.slides):
            print(f"  ⚠ Слайд {slide_idx} не существует — пропускаю")
            continue
        slide = prs.slides[slide_idx]
        if shape_idx >= len(slide.shapes):
            print(f"  ⚠ Shape {shape_idx} на слайде {slide_idx} не существует — пропускаю")
            continue
        shape = slide.shapes[shape_idx]
        if not shape.has_text_frame:
            print(f"  ⚠ Shape {shape_idx} на слайде {slide_idx} не имеет text_frame — пропускаю")
            continue

        tf = shape.text_frame
        # Очищаем первый параграф, первый run и вставляем плейсхолдер
        if tf.paragraphs and tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].text = placeholder
        elif tf.paragraphs:
            from pptx.oxml.ns import qn
            from lxml import etree
            para = tf.paragraphs[0]._p
            r = etree.SubElement(para, qn('a:r'))
            rPr = etree.SubElement(r, qn('a:rPr'), attrib={'lang': 'ru-RU'})
            t = etree.SubElement(r, qn('a:t'))
            t.text = placeholder
        print(f"  ✓ [{slide_idx},{shape_idx}] → {placeholder!r}")

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"\n✅ Шаблон сохранён: {output}")


def main() -> None:
    if "--inspect" in sys.argv:
        src = Path(sys.argv[sys.argv.index("--inspect") + 1]) if len(sys.argv) > 2 else SOURCE_PPTX
        if not src.exists():
            print(f"Файл не найден: {src}")
            sys.exit(1)
        inspect_pptx(src)
        return

    if not SOURCE_PPTX.exists():
        print(f"Исходный шаблон не найден: {SOURCE_PPTX}")
        print("Укажи путь вручную в переменной SOURCE_PPTX или скопируй файл в корень проекта.")
        sys.exit(1)

    print(f"Источник: {SOURCE_PPTX}")
    print(f"Вывод:    {OUTPUT_PPTX}\n")
    apply_placeholders(SOURCE_PPTX, OUTPUT_PPTX)


if __name__ == "__main__":
    main()
