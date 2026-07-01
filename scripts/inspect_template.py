"""Быстрая инспекция геометрии PPTX-шаблона (слайды, shape-ы, координаты)."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PATH = Path(__file__).parent.parent / "app" / "templates" / "report_template.pptx"


def main() -> None:
    if not PATH.exists():
        print(f"Не найден: {PATH}")
        return
    prs = Presentation(str(PATH))
    print(f"slides: {len(prs.slides)}")
    print(f"slide_size EMU: {prs.slide_width} x {prs.slide_height}")
    print(f"slide_size cm:  {Emu(prs.slide_width).cm:.1f} x {Emu(prs.slide_height).cm:.1f}")
    for i, slide in enumerate(prs.slides):
        print(f"\n=== slide {i} ({len(slide.shapes)} shapes) ===")
        for j, sh in enumerate(slide.shapes):
            txt = sh.text_frame.text[:60].replace("\n", " / ") if sh.has_text_frame else None
            top = f"{Emu(sh.top).cm:.1f}" if sh.top is not None else "?"
            left = f"{Emu(sh.left).cm:.1f}" if sh.left is not None else "?"
            w = f"{Emu(sh.width).cm:.1f}" if sh.width is not None else "?"
            h = f"{Emu(sh.height).cm:.1f}" if sh.height is not None else "?"
            print(f"  [{j}] {sh.shape_type} name={sh.name!r} "
                  f"top={top} left={left} w={w} h={h} text={txt!r}")


if __name__ == "__main__":
    main()
