"""
DOCX-сборщик отчёта «с нуля» (+ конвертация в PDF через LibreOffice).

Логика (согласовано с заказчиком, 2026-07-19):

Документ A4-портрет собирается программно из графических ассетов
(папка app/templates/new_assets) и текста. Разделителей между блоками больше
НЕТ — вместо них картинки-заголовки блоков во всю ширину страницы, а под каждой
идёт текст:

  • logo_top.png                — лого сверху во всю ширину
  • Профиль сотрудника          — смена / даты / департамент (в цвет) / ФИО
  • legend_block.png            → текст: контекст (легенда) смены
  • teachers_block.png          → текст: «преподский» блок (LLM)
  • tutors_block.png            → текст: «вожатский» блок (LLM)
  • logo_bottom.png             — лого снизу во всю ширину

Шрифты Calleo (calleo-regular.otf / calleo-semibold.otf) ВШИВАЮТСЯ внутрь .docx
(обфусцированный формат odttf по правилам OOXML), чтобы гарнитура отображалась
у того, кто открывает файл, даже без установки шрифта в системе.

ВАЖНО: этот модуль также экспортирует набор хелперов
(_hex_to_rgb, _resolve_shift_dates, get_department_color, get_department_name,
parse_numbered_answers, safe_pptx_filename), которые импортирует legacy
app/services/pptx_service.py — их сигнатуры НЕ меняем.
"""
from __future__ import annotations

import asyncio
import logging
import os
import re
import shutil
import subprocess
import zipfile

from datetime import date
from pathlib import Path

from docx import Document
from docx.enum.table import WD_ALIGN_VERTICAL, WD_ROW_HEIGHT_RULE
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls, qn
from docx.shared import Emu, Inches, Pt, RGBColor as DocxRGBColor







# RGBColor из python-pptx нужен для _hex_to_rgb (его импортирует pptx_service).
from pptx.dml.color import RGBColor as PptxRGBColor

from app.config import settings

logger = logging.getLogger(__name__)

# ─── Пути к шаблонам / ассетам ────────────────────────────────────────────────
TEMPLATE_DIR = Path(__file__).parent.parent / "templates"
ASSETS_DIR = TEMPLATE_DIR / "new_assets"

IMG_LOGO_TOP = ASSETS_DIR / "logo_top.png"
IMG_LOGO_BOTTOM = ASSETS_DIR / "logo_bottom.png"
IMG_LEGEND_BLOCK = ASSETS_DIR / "legend_block.png"
IMG_TEACHERS_BLOCK = ASSETS_DIR / "teachers_block.png"
IMG_TUTORS_BLOCK = ASSETS_DIR / "tutors_block.png"
# Вертикальная линия-разделитель ВНУТРИ профиля (между колонкой меток
# «Название смены / Дата смены / Департамент / ФИО» и колонкой значений).
IMG_PROFILE_V_SEP = ASSETS_DIR / "profile_vertical_separator.png"



FONT_REGULAR = ASSETS_DIR / "calleo-regular.otf"
FONT_SEMIBOLD = ASSETS_DIR / "calleo-semibold.otf"

FONT_NAME = "Calleo"

# ─── Размеры шрифта (pt) ──────────────────────────────────────────────────────
SIZE_BODY = 11
SIZE_PROFILE = 11
SIZE_PROFILE_LABEL = 11

# ─── Цвета ────────────────────────────────────────────────────────────────────
DARK_HEX = "0F1115"
GREY_HEX = "555555"

# ─── Цвета департаментов (hex без #) ──────────────────────────────────────────
# Официальные фирменные цвета департаментов (утверждено заказчиком).
# Должны совпадать с DEPARTMENTS в app/database/models.py.
DEPARTMENT_COLORS: dict[int, str] = {
    1: "F9423A",  # Департамент управления          — красный
    2: "FF672D",  # Департамент общественных связей — оранжево-красный
    3: "EDC731",  # Инженерный департамент          — жёлтый
    4: "242424",  # Департамент Икс                 — тёмно-серый / чёрный
    5: "50C787",  # Научный департамент             — зелёный
    6: "5A88FF",  # IT-департамент                  — синий
    7: "C061F3",  # Департамент дизайна             — фиолетовый
    8: "91D744",  # Проект 11                       — салатовый
    9: "FB4724",  # Летово Джун                     — оранжевый
}

DEPARTMENT_NAMES: dict[int, str] = {
    1: "Департамент управления",
    2: "Департамент общественных связей",
    3: "Инженерный департамент",
    4: "Департамент Икс",
    5: "Научный департамент",
    6: "IT-департамент",
    7: "Департамент дизайна",
    8: "Проект 11",
    9: "Летово Джун",
}

DEFAULT_COLOR = "E84130"  # fallback — фирменный красный «Летово»


def get_department_color(department_id: int) -> str:
    """Возвращает HEX-цвет (без #) для департамента."""
    return DEPARTMENT_COLORS.get(department_id, DEFAULT_COLOR)


def get_department_name(department_id: int) -> str:
    """Возвращает название департамента по ID."""
    return DEPARTMENT_NAMES.get(department_id, f"Департамент {department_id}")


# ─── Маркеры блоков в тексте отчёта ───────────────────────────────────────────
# LLM генерирует отчёт из двух блоков, разделяя их этими маркерами. Парсер ниже
# восстанавливает по ним «преподский» и «вожатский» тексты для документа.
TEACHER_MARKER = "=== ПРЕПОДСКИЙ БЛОК ==="
TUTOR_MARKER = "=== ВОЖАТСКИЙ БЛОК ==="


def split_two_blocks(text: str) -> tuple[str, str]:
    """Разбивает текст отчёта на (преподский_блок, вожатский_блок).

    Если маркеры не найдены — весь текст считаем преподским блоком, вожатский
    остаётся пустым (устойчиво к отклонениям LLM).
    """
    if not text:
        return "", ""
    t_idx = text.find(TEACHER_MARKER)
    v_idx = text.find(TUTOR_MARKER)

    if t_idx != -1 and v_idx != -1:
        teacher = text[t_idx + len(TEACHER_MARKER):v_idx].strip()
        tutor = text[v_idx + len(TUTOR_MARKER):].strip()
        return teacher, tutor
    if v_idx != -1:
        # Есть только вожатский маркер — до него преподский, после — вожатский.
        teacher = text[:v_idx].replace(TEACHER_MARKER, "").strip()
        tutor = text[v_idx + len(TUTOR_MARKER):].strip()
        return teacher, tutor
    if t_idx != -1:
        return text[t_idx + len(TEACHER_MARKER):].strip(), ""
    return text.strip(), ""


# ─── Утилиты имён файлов / дат ────────────────────────────────────────────────

def transliterate(text: str) -> str:
    table = {
        'а': 'a', 'б': 'b', 'в': 'v', 'г': 'g', 'д': 'd', 'е': 'e',
        'ё': 'yo', 'ж': 'zh', 'з': 'z', 'и': 'i', 'й': 'y', 'к': 'k',
        'л': 'l', 'м': 'm', 'н': 'n', 'о': 'o', 'п': 'p', 'р': 'r',
        'с': 's', 'т': 't', 'у': 'u', 'ф': 'f', 'х': 'kh', 'ц': 'ts',
        'ч': 'ch', 'ш': 'sh', 'щ': 'sch', 'ъ': '', 'ы': 'y', 'ь': '',
        'э': 'e', 'ю': 'yu', 'я': 'ya',
    }
    result = []
    for ch in text.lower():
        result.append(table.get(ch, ch))
    clean = ''.join(result)
    clean = re.sub(r'[^a-z0-9_]', '_', clean)
    return clean.strip('_')


# Символы, недопустимые в именах файлов (Windows/Linux). Кириллицу сохраняем.
_ILLEGAL_FILENAME_RE = re.compile(r'[\\/:*?"<>|\r\n\t]+')


def _clean_name_part(full_name: str) -> str:
    """«Фамилия Имя Отчество» → «Фамилия_Имя_Отчество» на русском.

    Пробелы заменяем на подчёркивания, вырезаем недопустимые в имени файла
    символы, схлопываем повторные подчёркивания. Кириллица остаётся как есть.
    """
    name = (full_name or "").strip()
    name = _ILLEGAL_FILENAME_RE.sub("", name)
    # Все пробельные последовательности → одиночное подчёркивание.
    name = re.sub(r"\s+", "_", name)
    name = re.sub(r"_+", "_", name).strip("_")
    return name or "Без_имени"


def safe_filename(student_name: str) -> Path:
    """Отчет_Сотрудника_Фамилия_Имя_Отчество.docx (на русском).

    Если отчества нет — «Отчет_Сотрудника_Фамилия_Имя.docx».
    """
    return Path(f"Отчет_Сотрудника_{_clean_name_part(student_name)}.docx")



def safe_pptx_filename(student_name: str) -> Path:
    """report_IvanovIvan.pptx (для legacy pptx_service)."""
    return Path(f"report_{transliterate(student_name)}.pptx")


def safe_archive_name(shift_name: str) -> str:
    """reports_Smena1IT_20260621.zip"""
    return f"reports_{transliterate(shift_name)}_{date.today().strftime('%Y%m%d')}.zip"


def _resolve_shift_dates(shift) -> str:
    """Собирает строку с датами смены из доступных полей модели."""
    dates = getattr(shift, "dates", None)
    if dates:
        return str(dates)
    start = getattr(shift, "start_date", None) or getattr(shift, "date_start", None)
    end = getattr(shift, "end_date", None) or getattr(shift, "date_end", None)
    try:
        if start and end:
            return f"{start:%d.%m} – {end:%d.%m.%Y}"
        if start:
            return f"{start:%d.%m.%Y}"
    except (ValueError, TypeError):
        pass
    return ""


def _strip_markdown_prefix(line: str) -> str:
    """Убирает markdown-обёртки в начале строки (** , * , #, >, -, •, пробелы)."""
    return re.sub(r"^[\s>#*•\-\u2013\u2014]+", "", line)


def parse_numbered_answers(text: str) -> dict[int, str]:
    """
    Fallback-парсер пронумерованных ответов «1. …», «2) …» из текста отчёта.
    Оставлен для обратной совместимости (используется legacy pptx_service.py).
    """
    if not text:
        return {}
    answers: dict[int, str] = {}
    current_num: int | None = None
    buf: list[str] = []
    num_re = re.compile(r"^(\d{1,2})\s*[.):\-\u2013\u2014]\s*(.*)$")
    num_only_re = re.compile(r"^(\d{1,2})\s*[.):\-\u2013\u2014]?\s*$")
    stop_re = re.compile(r"(итогов\w*\s+отч|={3,})", re.IGNORECASE)

    def _flush() -> None:
        if current_num is not None:
            answers[current_num] = " ".join(buf).strip()

    for raw_line in text.splitlines():
        if stop_re.search(raw_line):
            break
        line = _strip_markdown_prefix(raw_line).strip()
        if not line:
            continue
        m_only = num_only_re.match(line)
        if m_only:
            _flush()
            current_num = int(m_only.group(1))
            buf = []
            continue
        m = num_re.match(line)
        if m:
            _flush()
            current_num = int(m.group(1))
            buf = [m.group(2).strip()] if m.group(2).strip() else []
        elif current_num is not None:
            buf.append(line)

    _flush()
    return {k: v for k, v in answers.items() if v}


# ─── Конвертация цветов ───────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str) -> PptxRGBColor:
    """'C0392B' → pptx RGBColor(192, 57, 43). Используется legacy pptx_service."""
    hex_color = hex_color.lstrip("#")
    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    return PptxRGBColor(r, g, b)


def _hex_to_docx_rgb(hex_color: str) -> DocxRGBColor:
    """'C0392B' → docx RGBColor(192, 57, 43)."""
    hex_color = hex_color.lstrip("#")
    return DocxRGBColor(
        int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    )


# ─── Геометрия изображений ────────────────────────────────────────────────────

def _png_size(path: Path) -> tuple[int, int]:
    """Возвращает (width, height) PNG в пикселях, читая заголовок IHDR.

    Без PIL: у PNG первые 8 байт — сигнатура, затем чанк IHDR, где ширина и
    высота лежат как два big-endian uint32 по смещениям 16 и 20.
    """
    with open(path, "rb") as f:
        head = f.read(24)
    if len(head) < 24 or head[:8] != b"\x89PNG\r\n\x1a\n":
        # Не PNG или битый заголовок — вернём квадрат, чтобы не делить на ноль.
        return 1, 1
    w = int.from_bytes(head[16:20], "big")
    h = int.from_bytes(head[20:24], "big")
    return (w or 1), (h or 1)


def _rendered_height_emu(path: Path, width_emu: int) -> int:
    """Высота изображения в EMU при отрисовке на заданную ширину (с сохранением
    пропорций). Нужна, чтобы зарезервировать место под нижний логотип в footer."""
    w_px, h_px = _png_size(path)
    return int(width_emu * h_px / w_px)


def _set_table_indent(table, indent_emu: int) -> None:
    """Задаёт левый отступ таблицы (w:tblInd) — так профиль отодвигается от края
    страницы, когда поля секции равны нулю."""
    tbl_pr = table._tbl.tblPr
    ind = tbl_pr.find(qn("w:tblInd"))
    if ind is None:
        ind = tbl_pr.makeelement(qn("w:tblInd"), {})
        tbl_pr.append(ind)
    ind.set(qn("w:w"), str(int(indent_emu / 635)))  # EMU → twips (1 twip = 635 EMU)
    ind.set(qn("w:type"), "dxa")



# ─── Вшивание шрифтов Calleo внутрь .docx ─────────────────────────────────────
# По правилам OOXML (ECMA-376 §17.8) встроенные шрифты Word хранятся в
# обфусцированном виде (.odttf): первые 32 байта файла XOR-ятся с 16-байтным
# ключом (GUID). GUID хранится в атрибуте w:fontKey внутри word/fontTable.xml.
# Всё обёрнуто в try/except: при любой ошибке .docx остаётся валидным (просто без
# встроенного шрифта).

_W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_FONT_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"
_FONTTABLE_REL_TYPE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/fontTable"

# Порядок дочерних элементов CT_Settings (нужен, чтобы Word не «чинил» файл).
# Достаточно ранней части последовательности; неизвестным тегам присваиваем
# большой индекс, чтобы embedTrueTypeFonts встал перед ними (они все позже).
_SETTINGS_ORDER = {
    "writeProtection": 0, "view": 1, "zoom": 2, "removePersonalInformation": 3,
    "doNotDisplayPageBoundaries": 4, "displayBackgroundShape": 5,
    "printPostScriptOverText": 6, "printFractionalCharacterWidth": 7,
    "printFormsData": 8, "embedTrueTypeFonts": 9, "embedSystemFonts": 10,
    "saveSubsetFonts": 11, "saveFormsData": 12, "mirrorMargins": 13,
    "alignBordersAndEdges": 14, "bordersDoNotSurroundHeader": 15,
    "bordersDoNotSurroundFooter": 16, "gutterAtTop": 17, "hideSpellingErrors": 18,
    "hideGrammaticalErrors": 19, "defaultTabStop": 33,
    "characterSpacingControl": 40, "compat": 50, "rsids": 51,
}


def _make_font_key() -> tuple[str, bytes]:
    """Возвращает (w:fontKey-строка «{GUID}», 16-байтный ключ для XOR).

    Строка формируется так, чтобы Word, распарсив GUID и вызвав внутреннее
    Guid.ToByteArray() (mixed-endian), получил ровно те же 16 байт, которыми мы
    обфусцируем шрифт.
    """
    kb = os.urandom(16)
    guid = (
        "{%02X%02X%02X%02X-%02X%02X-%02X%02X-%02X%02X-%02X%02X%02X%02X%02X%02X}"
        % (
            kb[3], kb[2], kb[1], kb[0],
            kb[5], kb[4],
            kb[7], kb[6],
            kb[8], kb[9],
            kb[10], kb[11], kb[12], kb[13], kb[14], kb[15],
        )
    )
    return guid, kb


def _obfuscate_font(font_bytes: bytes, kb: bytes) -> bytes:
    """XOR-обфускация первых 32 байт шрифта ключом kb (по OOXML)."""
    data = bytearray(font_bytes)
    limit = min(32, len(data))
    for i in range(limit):
        data[i] ^= kb[i % 16]
    return bytes(data)


def _next_rid(existing: set) -> str:
    i = 1
    while f"rId{i}" in existing:
        i += 1
    return f"rId{i}"


def _insert_settings_ordered(settings_root, new_el, ns: str) -> None:
    """Вставляет new_el в settings.xml, сохраняя порядок схемы CT_Settings."""
    tag = new_el.tag.split("}")[-1]
    new_idx = _SETTINGS_ORDER.get(tag, 999)
    for child in settings_root:
        child_tag = child.tag.split("}")[-1]
        child_idx = _SETTINGS_ORDER.get(child_tag, 999)
        if child_idx > new_idx:
            child.addprevious(new_el)
            return
    settings_root.append(new_el)


def _embed_calleo_fonts_docx(docx_path: Path) -> None:
    """Встраивает шрифты Calleo (regular + semibold) в .docx как odttf."""
    if not FONT_REGULAR.exists():
        logger.warning("Font asset missing, skip embed: %s", FONT_REGULAR)
        return
    try:
        from lxml import etree

        with zipfile.ZipFile(docx_path, "r") as zin:
            data = {n: zin.read(n) for n in zin.namelist()}

        # --- 1. Обфусцированные данные шрифтов ---
        reg_guid, reg_kb = _make_font_key()
        bold_guid, bold_kb = _make_font_key()
        semi_src = FONT_SEMIBOLD if FONT_SEMIBOLD.exists() else FONT_REGULAR
        data["word/fonts/font1.odttf"] = _obfuscate_font(FONT_REGULAR.read_bytes(), reg_kb)
        data["word/fonts/font2.odttf"] = _obfuscate_font(semi_src.read_bytes(), bold_kb)

        # --- 2. word/fontTable.xml ---
        fonts_el = etree.Element(f"{{{_W}}}fonts", nsmap={"w": _W, "r": _R})
        font_el = etree.SubElement(fonts_el, f"{{{_W}}}font")
        font_el.set(f"{{{_W}}}name", FONT_NAME)
        er = etree.SubElement(font_el, f"{{{_W}}}embedRegular")
        er.set(f"{{{_R}}}id", "rIdFontReg")
        er.set(f"{{{_W}}}fontKey", reg_guid)
        eb = etree.SubElement(font_el, f"{{{_W}}}embedBold")
        eb.set(f"{{{_R}}}id", "rIdFontBold")
        eb.set(f"{{{_W}}}fontKey", bold_guid)
        data["word/fontTable.xml"] = etree.tostring(
            fonts_el, xml_declaration=True, encoding="UTF-8", standalone=True
        )

        # --- 3. word/_rels/fontTable.xml.rels ---
        rels_root = etree.Element(f"{{{_REL}}}Relationships")
        for rid, target in (
            ("rIdFontReg", "fonts/font1.odttf"),
            ("rIdFontBold", "fonts/font2.odttf"),
        ):
            rel = etree.SubElement(rels_root, f"{{{_REL}}}Relationship")
            rel.set("Id", rid)
            rel.set("Type", _FONT_REL_TYPE)
            rel.set("Target", target)
        data["word/_rels/fontTable.xml.rels"] = etree.tostring(
            rels_root, xml_declaration=True, encoding="UTF-8", standalone=True
        )

        # --- 4. word/_rels/document.xml.rels: связь на fontTable.xml ---
        doc_rels_name = "word/_rels/document.xml.rels"
        doc_rels = etree.fromstring(data[doc_rels_name])
        has_fonttable = any(
            r.get("Type") == _FONTTABLE_REL_TYPE for r in doc_rels
        )
        if not has_fonttable:
            existing_ids = {r.get("Id") for r in doc_rels}
            rel = etree.SubElement(doc_rels, f"{{{_REL}}}Relationship")
            rel.set("Id", _next_rid(existing_ids))
            rel.set("Type", _FONTTABLE_REL_TYPE)
            rel.set("Target", "fontTable.xml")
            data[doc_rels_name] = etree.tostring(
                doc_rels, xml_declaration=True, encoding="UTF-8", standalone=True
            )

        # --- 5. word/settings.xml: <w:embedTrueTypeFonts/> ---
        settings_name = "word/settings.xml"
        settings_root = etree.fromstring(data[settings_name])
        if settings_root.find(f"{{{_W}}}embedTrueTypeFonts") is None:
            ett = etree.Element(f"{{{_W}}}embedTrueTypeFonts")
            _insert_settings_ordered(settings_root, ett, _W)
            data[settings_name] = etree.tostring(
                settings_root, xml_declaration=True, encoding="UTF-8", standalone=True
            )

        # --- 6. [Content_Types].xml: Default odttf + Override fontTable ---
        ct_name = "[Content_Types].xml"
        ct_root = etree.fromstring(data[ct_name])
        if not any(
            d.get("Extension") == "odttf" for d in ct_root.findall(f"{{{_CT}}}Default")
        ):
            d = etree.SubElement(ct_root, f"{{{_CT}}}Default")
            d.set("Extension", "odttf")
            d.set("ContentType", "application/vnd.openxmlformats-officedocument.obfuscatedFont")
        if not any(
            o.get("PartName") == "/word/fontTable.xml"
            for o in ct_root.findall(f"{{{_CT}}}Override")
        ):
            o = etree.SubElement(ct_root, f"{{{_CT}}}Override")
            o.set("PartName", "/word/fontTable.xml")
            o.set(
                "ContentType",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.fontTable+xml",
            )
        data[ct_name] = etree.tostring(
            ct_root, xml_declaration=True, encoding="UTF-8", standalone=True
        )

        # --- 7. Перезаписываем zip ---
        tmp_path = docx_path.with_suffix(".tmp.docx")
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for name, payload in data.items():
                zout.writestr(name, payload)
        shutil.move(str(tmp_path), str(docx_path))
        logger.info("Calleo fonts embedded into %s", docx_path.name)

    except Exception as e:  # noqa: BLE001 — не роняем экспорт из-за встраивания
        logger.warning("DOCX font embedding skipped (%s): %s", docx_path.name, e)


# ─── DOCX ─────────────────────────────────────────────────────────────────────

class DocxService:
    def __init__(self) -> None:
        self.reports_dir = Path(settings.reports_dir)
        self.reports_dir.mkdir(parents=True, exist_ok=True)

    # ── низкоуровневые помощники вёрстки ──────────────────────────────────────
    @staticmethod
    def _apply_font(run, size: int, bold: bool, color_hex: str) -> None:
        run.font.name = FONT_NAME
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _hex_to_docx_rgb(color_hex)
        # Явно проставляем гарнитуру для всех диапазонов (в т.ч. кириллицы).
        rpr = run._element.get_or_add_rPr()
        rfonts = rpr.find(qn("w:rFonts"))
        if rfonts is None:
            rfonts = rpr.makeelement(qn("w:rFonts"), {})
            rpr.insert(0, rfonts)
        for attr in ("w:ascii", "w:hAnsi", "w:cs", "w:eastAsia"):
            rfonts.set(qn(attr), FONT_NAME)

    def _add_body(self, doc, text: str, side_indent) -> None:
        """Добавляет абзацы тела блока (разбивает по переносам строк).

        side_indent — отступ слева/справа: поля страницы обнулены (чтобы
        картинки были во всю ширину), поэтому текст отодвигаем от края вручную.
        """
        text = (text or "").strip() or "—"
        # Абзацы разделяем по пустым строкам; одиночные \n оставляем как разрыв.
        chunks = [c.strip() for c in re.split(r"\n\s*\n", text) if c.strip()]
        if not chunks:
            chunks = [text]
        for chunk in chunks:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
            pf = p.paragraph_format
            pf.space_before = Pt(2)
            pf.space_after = Pt(6)
            pf.left_indent = side_indent
            pf.right_indent = side_indent
            # Красная строка: каждый абзац начинается с отступа первой строки.
            pf.first_line_indent = Inches(0.5)
            lines = chunk.split("\n")

            for i, line in enumerate(lines):
                run = p.add_run(line)
                self._apply_font(run, SIZE_BODY, False, DARK_HEX)
                if i < len(lines) - 1:
                    run.add_break()

    def _add_fullwidth_image(self, doc, path: Path, page_width,
                             space_before: int = 6, space_after: int = 2,
                             keep_with_next: bool = False) -> None:
        """Вставляет изображение во ВСЮ ширину страницы (края в край).

        Ширина = полная ширина страницы (поля секции обнулены), высота
        подбирается автоматически с сохранением пропорций — перспектива не
        ломается, т.к. задаём только width.

        keep_with_next=True «сцепляет» картинку-заголовок блока со следующим
        абзацем: если перед ней много текста и заголовок не помещается внизу
        страницы, он вместе с началом текста переносится на новую страницу
        (картинка не «прилипает» к нижнему краю и не обрезается).
        """
        if not path.exists():
            logger.warning("Asset missing: %s", path)
            return
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_before = Pt(space_before)
        p.paragraph_format.space_after = Pt(space_after)
        # Абзац-контейнер картинки тоже без отступов, иначе появится сдвиг вправо.
        p.paragraph_format.left_indent = Emu(0)
        p.paragraph_format.right_indent = Emu(0)
        if keep_with_next:
            p.paragraph_format.keep_with_next = True
            p.paragraph_format.keep_together = True
        run = p.add_run()
        run.add_picture(str(path), width=page_width)


    def _add_profile(self, doc, shift_name: str, shift_dates: str,
                     dep_name: str, dep_color_hex: str, student_name: str,
                     content_width, side_indent) -> None:
        """Профиль сотрудника: три колонки (метка | вертикальный разделитель |
        значение), без рамок.

        Строки:
            Название смены │ …
            Дата смены     │ …
            Департамент    │ …
            ФИО            │ …

        В средней узкой колонке в каждой строке стоит вертикальная линия-
        разделитель profile_vertical_separator.png — она визуально делит метки
        и значения на всю высоту профиля.

        content_width — ширина текстовой области (страница минус боковые
        отступы), side_indent — левый отступ от края страницы (поля обнулены).
        """
        rows = [
            ("Название смены", (shift_name or "—").upper(), DARK_HEX),
            ("Дата смены", (shift_dates or "—").upper(), DARK_HEX),
            ("Департамент", (dep_name or "—").upper(), dep_color_hex),
            ("ФИО", (student_name or "—").upper(), DARK_HEX),
        ]
        table = doc.add_table(rows=len(rows), cols=3)
        table.autofit = False
        _set_table_indent(table, int(side_indent))

        # Ширины колонок: метки | разделитель | значения.
        sep_w = Inches(0.14)
        label_w = Inches(1.9)
        value_w = content_width - label_w - sep_w
        col_widths = (label_w, sep_w, value_w)

        # Фиксированная высота каждой строки — чтобы точно рассчитать высоту
        # ОДНОЙ вертикальной линии-разделителя на весь профиль.
        row_h_pt = 24
        for row in table.rows:
            row.height = Pt(row_h_pt)
            row.height_rule = WD_ROW_HEIGHT_RULE.EXACTLY

        # Метки (колонка 0) и значения (колонка 2).
        for r_idx, (label, value, color_hex) in enumerate(rows):
            lc = table.cell(r_idx, 0)
            vc = table.cell(r_idx, 2)
            for cell, w in ((lc, label_w), (vc, value_w)):
                cell.width = w
                cell.vertical_alignment = WD_ALIGN_VERTICAL.CENTER

            lp = lc.paragraphs[0]
            lp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
            lrun = lp.add_run(label)
            self._apply_font(lrun, SIZE_PROFILE_LABEL, False, GREY_HEX)

            vp = vc.paragraphs[0]
            vp.alignment = WD_ALIGN_PARAGRAPH.LEFT
            # Отступ значения от вертикальной линии-разделителя. Линия теперь
            # плавающая (перед текстом), поэтому paragraph-отступ работает
            # корректно: текст отходит правее, а линия остаётся на месте.
            vp.paragraph_format.left_indent = Inches(0.22)
            vrun = vp.add_run(value)
            self._apply_font(vrun, SIZE_PROFILE, True, color_hex)




        # Средняя колонка: объединяем все ячейки в ОДНУ и кладём туда ОДНУ
        # вертикальную линию-разделитель на всю высоту профиля — как ПЛАВАЮЩИЙ
        # объект «ПЕРЕД текстом» (in front of text). Линия лежит ПОВЕРХ всего,
        # поэтому текст значений её не «подтёсывает» и не наезжает на неё.
        merged = table.cell(0, 1)
        for r_idx in range(1, len(rows)):
            merged = merged.merge(table.cell(r_idx, 1))
        merged.width = sep_w
        merged.vertical_alignment = WD_ALIGN_VERTICAL.CENTER
        sp = merged.paragraphs[0]
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sp.paragraph_format.space_before = Pt(0)
        sp.paragraph_format.space_after = Pt(0)
        if IMG_PROFILE_V_SEP.exists():
            srun = sp.add_run()
            # Высота линии = суммарная высота строк профиля. Ширину не задаём —
            # сохраняем пропорции тонкой линии.
            srun.add_picture(str(IMG_PROFILE_V_SEP), height=Pt(row_h_pt * len(rows)))
            # Переводим линию в режим «перед текстом» (лежит поверх содержимого).
            self._float_in_front_of_text(srun)
        else:
            logger.warning("Asset missing: %s", IMG_PROFILE_V_SEP)



    def _add_bottom_logo(self, doc, page_width) -> None:
        """Нижнее лого — во всю ширину страницы, ПРИЖАТО к нижнему краю
        последней страницы (не сразу после текста). НЕ колонтитул.

        Реализовано плавающим (floating) изображением, привязанным по вертикали
        к низу СТРАНИЦЫ (relativeFrom="page", align="bottom") и по центру по
        горизонтали. Такой объект «садится» в самый низ той страницы, где стоит
        его абзац-якорь (последний в документе), независимо от количества текста.
        """
        if not IMG_LOGO_BOTTOM.exists():
            logger.warning("Asset missing: %s", IMG_LOGO_BOTTOM)
            return

        from lxml import etree

        width_emu = int(page_width)
        height_emu = _rendered_height_emu(IMG_LOGO_BOTTOM, width_emu)

        # Абзац-якорь: к нему привязан плавающий объект. Сам абзац пустой.
        p = doc.add_paragraph()
        p.paragraph_format.space_before = Pt(0)
        p.paragraph_format.space_after = Pt(0)
        p.paragraph_format.left_indent = Emu(0)
        p.paragraph_format.right_indent = Emu(0)
        run = p.add_run()
        run.add_picture(str(IMG_LOGO_BOTTOM), width=page_width)

        # Достаём inline-рисунок и превращаем его в плавающий anchor.
        drawing = run._element.find(qn("w:drawing"))
        inline = drawing.find(qn("wp:inline"))
        graphic = inline.find(qn("a:graphic"))
        docpr = inline.find(qn("wp:docPr"))
        doc_id = docpr.get("id") if docpr is not None else "999"
        doc_name = docpr.get("name") if docpr is not None else "logo_bottom"
        graphic_xml = etree.tostring(graphic).decode()

        anchor_xml = (
            f'<wp:anchor {nsdecls("wp", "a", "r", "pic")} '
            'behindDoc="0" distT="0" distB="0" distL="0" distR="0" '
            'simplePos="0" locked="0" layoutInCell="1" allowOverlap="1" '
            'relativeHeight="0">'
            '<wp:simplePos x="0" y="0"/>'
            '<wp:positionH relativeFrom="page"><wp:align>center</wp:align></wp:positionH>'
            '<wp:positionV relativeFrom="page"><wp:align>bottom</wp:align></wp:positionV>'
            f'<wp:extent cx="{width_emu}" cy="{height_emu}"/>'
            '<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            '<wp:wrapNone/>'
            f'<wp:docPr id="{doc_id}" name="{doc_name}"/>'
            '<wp:cNvGraphicFramePr/>'
            f'{graphic_xml}'
            '</wp:anchor>'
        )
        anchor = parse_xml(anchor_xml)
        drawing.replace(inline, anchor)




    def _add_page_break(self, doc) -> None:
        """Явный разрыв страницы: следующий контент — с новой страницы."""
        doc.add_page_break()

    def _float_in_front_of_text(self, run) -> None:
        """Превращает инлайновую картинку в плавающий объект «ПЕРЕД текстом».

        Аналог опции Word «Обтекание текстом → Перед текстом» (In Front of Text):
        картинка лежит ПОВЕРХ содержимого (behindDoc="0", <wp:wrapNone/>), поэтому
        текст рядом её не «подтёсывает» и не наезжает. Позиция берётся из текущего
        места вставки (по горизонтали — от колонки, по вертикали — от абзаца).
        """
        from lxml import etree

        drawing = run._element.find(qn("w:drawing"))
        if drawing is None:
            return
        inline = drawing.find(qn("wp:inline"))
        if inline is None:
            return
        extent = inline.find(qn("wp:extent"))
        cx = extent.get("cx") if extent is not None else "0"
        cy = extent.get("cy") if extent is not None else "0"
        graphic = inline.find(qn("a:graphic"))
        docpr = inline.find(qn("wp:docPr"))
        doc_id = docpr.get("id") if docpr is not None else "998"
        doc_name = docpr.get("name") if docpr is not None else "v_sep"
        graphic_xml = etree.tostring(graphic).decode()

        anchor_xml = (
            f'<wp:anchor {nsdecls("wp", "a", "r", "pic")} '
            'behindDoc="0" distT="0" distB="0" distL="0" distR="0" '
            'simplePos="0" locked="0" layoutInCell="1" allowOverlap="1" '
            'relativeHeight="251658240">'
            '<wp:simplePos x="0" y="0"/>'
            '<wp:positionH relativeFrom="column">'
            '<wp:posOffset>0</wp:posOffset></wp:positionH>'
            '<wp:positionV relativeFrom="paragraph">'
            '<wp:posOffset>0</wp:posOffset></wp:positionV>'
            f'<wp:extent cx="{cx}" cy="{cy}"/>'
            '<wp:effectExtent l="0" t="0" r="0" b="0"/>'
            '<wp:wrapNone/>'
            f'<wp:docPr id="{doc_id}" name="{doc_name}"/>'
            '<wp:cNvGraphicFramePr/>'
            f'{graphic_xml}'
            '</wp:anchor>'
        )
        anchor = parse_xml(anchor_xml)
        drawing.replace(inline, anchor)


    # ── основной метод генерации ──────────────────────────────────────────────
    def generate(

        self,
        report,
        student,
        shift,
        teacher,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        """Собирает DOCX и возвращает путь к файлу."""
        # Департамент берётся из РЕБЁНКА (student.department_number проставляется
        # в хендлере), фолбэк — на смену.
        dep_number = getattr(student, "department_number", None)
        if dep_number is None:
            dep_number = getattr(shift, "department_id", None) or 0
        dep_color_hex = get_department_color(dep_number)
        dep_name = get_department_name(dep_number)

        teacher_block, tutor_block = split_two_blocks(report.generated_text or "")
        legend_text = (shift_context or "").strip() or "Контекст смены не указан."

        doc = Document()
        section = doc.sections[0]
        section.page_width = Inches(8.27)     # A4 портрет
        section.page_height = Inches(11.69)
        # Поля страницы обнулены по бокам и сверху: так лого и картинки-заголовки
        # блоков занимают ВСЮ ширину страницы (края в край), а верхнее лого не
        # имеет отступа от начала страницы. Нижнее поле выставит футер под лого.
        section.left_margin = Emu(0)
        section.right_margin = Emu(0)
        section.top_margin = Emu(0)
        section.bottom_margin = Emu(0)
        # Полная ширина страницы — для изображений «во всю ширину».
        page_width = section.page_width
        # Боковой отступ для ТЕКСТА/профиля (картинки идут в край, текст — с полями).
        side_indent = Inches(0.6)

        # Лого сверху — во всю ширину, строго в начале страницы (без отступа).
        self._add_fullwidth_image(
            doc, IMG_LOGO_TOP, page_width, space_before=0, space_after=4
        )

        # Профиль сотрудника
        self._add_profile(
            doc,
            shift_name=shift.name or "",
            shift_dates=_resolve_shift_dates(shift),
            dep_name=dep_name,
            dep_color_hex=dep_color_hex,
            student_name=student.full_name or "",
            content_width=page_width - side_indent * 2,
            side_indent=side_indent,
        )

        # Блок 1. Легенда смены — СТРОГО на первой странице.
        self._add_fullwidth_image(doc, IMG_LEGEND_BLOCK, page_width)
        self._add_body(doc, legend_text, side_indent)

        # Блок 2. Преподский — начинается с НОВОЙ (второй) страницы.
        self._add_page_break(doc)
        self._add_fullwidth_image(doc, IMG_TEACHERS_BLOCK, page_width)
        self._add_body(doc, teacher_block, side_indent)

        # Блок 3. Вожатский — идёт СРАЗУ после текста преподского блока
        # (без разрыва страницы). keep_with_next: если преподский текст длинный
        # и картинка-заголовок вожатского блока не помещается внизу страницы,
        # она НЕ обрезается краем, а целиком переносится на новую страницу.
        self._add_fullwidth_image(
            doc, IMG_TUTORS_BLOCK, page_width, keep_with_next=True
        )
        self._add_body(doc, tutor_block, side_indent)




        # Лого снизу — во всю ширину, без отступов, в конце тела документа
        # (на последней странице). НЕ колонтитул.
        self._add_bottom_logo(doc, page_width)


        output_path = self.reports_dir / safe_filename(student.full_name)

        doc.save(str(output_path))
        _embed_calleo_fonts_docx(output_path)
        logger.info("DOCX built: %s for student=%r", output_path, student.full_name)
        return output_path

    # ── конвертация в PDF (через LibreOffice/soffice) ─────────────────────────
    def _to_pdf(self, docx_path: Path) -> Path:
        """Конвертирует .docx → .pdf через headless LibreOffice."""
        outdir = docx_path.parent
        subprocess.run(
            [
                "soffice", "--headless", "--nologo", "--nofirststartwizard",
                "--convert-to", "pdf", "--outdir", str(outdir), str(docx_path),
            ],
            check=True,
            timeout=180,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        pdf_path = docx_path.with_suffix(".pdf")
        if not pdf_path.exists():
            raise RuntimeError(f"PDF не создан: {pdf_path}")
        logger.info("PDF built: %s", pdf_path.name)
        return pdf_path

    def generate_pdf(
        self,
        report,
        student,
        shift,
        teacher,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        """Собирает DOCX и конвертирует его в PDF, возвращает путь к PDF."""
        docx_path = self.generate(
            report=report, student=student, shift=shift, teacher=teacher,
            q_answers=q_answers, shift_context=shift_context,
        )
        return self._to_pdf(docx_path)

    # ── асинхронные обёртки (тяжёлая работа — в отдельном потоке) ──────────────
    async def generate_async(
        self,
        report,
        student,
        shift,
        teacher,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        return await asyncio.to_thread(
            self.generate,
            report=report, student=student, shift=shift, teacher=teacher,
            q_answers=q_answers, shift_context=shift_context,
        )

    async def generate_pdf_async(
        self,
        report,
        student,
        shift,
        teacher,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        return await asyncio.to_thread(
            self.generate_pdf,
            report=report, student=student, shift=shift, teacher=teacher,
            q_answers=q_answers, shift_context=shift_context,
        )
