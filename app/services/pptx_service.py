"""
PPTX-сборщик отчёта «с нуля».

В отличие от старого подхода (мутация report_template.pptx с плейсхолдерами и
ручным сдвигом фигур), здесь презентация СОБИРАЕТСЯ программно из графических
ассетов (папка app/templates/assets) и текста. Это устраняет «разъезжающуюся»
вёрстку и пропадающие поля.

Структура документа (согласовано с заказчиком):

  Страница 1 (отдельная):
    • logo_top.png                    — лого сверху по центру
    • separator_corner.png            — заглавный угловой разделитель
    • Профиль сотрудника (2 колонки + profile_vertical_separator.png):
        Название смены  | shift_name
        Дата смены      | shift_dates
        Департамент     | department_name (в цвет департамента)
        ФИО             | student_name
    • separator_line.png
    • ЛЕГЕНДА СМЕНЫ + текст (shift_context)

  Страницы 2..N (авто-пагинация):
    • 5 блоков с вопросами; после каждого блока separator_line.png,
      кроме последнего блока.

  Последняя страница (отдельная):
    • first_separator_for_last_page.png  — заглавный разделитель сверху
    • ОБЩЕЕ ЗАКЛЮЧЕНИЕ + текст (итоговый отчёт)
    • last_separator.png
    • logo_bottom.png                    — красивое лого снизу по центру

  arrow.png — в правом нижнем углу каждой страницы, КРОМЕ последней.

Шрифты (ассеты calleo-regular.otf / calleo-semibold.otf лежат рядом — их нужно
установить в системе того, кто открывает файл, чтобы гарнитура отобразилась):
  • Названия блоков — Calleo, полужирный, 14 pt (в цвет департамента)
  • Подзаголовки вопросов — Calleo, обычный, 14 pt
  • Остальной текст — Calleo, обычный, 10 pt
"""

from __future__ import annotations

import logging
import re
import shutil
import subprocess
import zipfile
from pathlib import Path


from pptx import Presentation

from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.text import MSO_AUTO_SIZE

from app.config import settings
from app.database.models import Report, Student, Shift, User
from app.services.docx_service import (
    _hex_to_rgb,
    _resolve_shift_dates,
    get_department_color,
    get_department_name,
    parse_numbered_answers,
    safe_pptx_filename,
)

logger = logging.getLogger(__name__)

# ─── Пути ─────────────────────────────────────────────────────────────────────
TEMPLATE_DIR = Path(__file__).parent.parent / "templates"
ASSETS_DIR = TEMPLATE_DIR / "assets"
BASE_TEMPLATE = TEMPLATE_DIR / "report_template.pptx"

# ─── Ассеты ───────────────────────────────────────────────────────────────────
IMG_LOGO_TOP = ASSETS_DIR / "logo_top.png"
IMG_LOGO_BOTTOM = ASSETS_DIR / "logo_bottom.png"
IMG_SEP_CORNER = ASSETS_DIR / "separator_corner.png"
IMG_SEP_LINE = ASSETS_DIR / "separator_line.png"
IMG_SEP_PROFILE_V = ASSETS_DIR / "profile_vertical_separator.png"
IMG_SEP_FIRST_LAST = ASSETS_DIR / "first_separator_for_last_page.png"
IMG_SEP_LAST = ASSETS_DIR / "last_separator.png"
IMG_ARROW = ASSETS_DIR / "arrow.png"

# ─── Шрифты / размеры ─────────────────────────────────────────────────────────
FONT = "Calleo"
SIZE_BLOCK = 14      # название блока (полужирный)
SIZE_QUESTION = 11   # текст вопроса (полужирный) — совпадает по размеру с ответом
SIZE_BODY = 11       # обычный текст / ответ педагога
SIZE_PROFILE = 11    # профиль


FONT_REGULAR = ASSETS_DIR / "calleo-regular.otf"
FONT_SEMIBOLD = ASSETS_DIR / "calleo-semibold.otf"

# XML-namespaces для правки presentation.xml
_NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
_NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
_REL_FONT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/font"


def _embed_calleo_fonts(pptx_path: Path) -> None:
    """
    Встраивает шрифты Calleo (regular + semibold) в .pptx, чтобы гарнитура
    отображалась даже если она не установлена у зрителя.

    Делается через постобработку zip-архива pptx:
      • ppt/fonts/font1.fntdata (regular), font2.fntdata (semibold)
      • [Content_Types].xml   — Default для расширения fntdata
      • ppt/presentation.xml  — embedTrueTypeFonts + <p:embeddedFontLst>
      • ppt/_rels/presentation.xml.rels — relationships на файлы шрифтов

    Всё обёрнуто в try/except: при любой ошибке исходный файл остаётся как есть
    (шрифт просто не встроится, но презентация не «сломается»).
    """
    if not FONT_REGULAR.exists():
        logger.warning("Font asset missing, skip embed: %s", FONT_REGULAR)
        return
    try:
        from lxml import etree

        with zipfile.ZipFile(pptx_path, "r") as zin:
            names = zin.namelist()
            data = {n: zin.read(n) for n in names}

        # --- 1. presentation.xml.rels: добавляем relationships на шрифты ---
        rels_name = "ppt/_rels/presentation.xml.rels"
        rels_root = etree.fromstring(data[rels_name])
        existing_ids = {
            r.get("Id") for r in rels_root
        }
        rid_font_reg = _next_rid(existing_ids)
        existing_ids.add(rid_font_reg)
        rid_font_semi = _next_rid(existing_ids)
        existing_ids.add(rid_font_semi)

        for rid, target in (
            (rid_font_reg, "fonts/font1.fntdata"),
            (rid_font_semi, "fonts/font2.fntdata"),
        ):
            rel = etree.SubElement(rels_root, f"{{{_NS_REL}}}Relationship")
            rel.set("Id", rid)
            rel.set("Type", _REL_FONT)
            rel.set("Target", target)
        data[rels_name] = etree.tostring(rels_root, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)

        # --- 2. presentation.xml: embedTrueTypeFonts + embeddedFontLst ---
        pres_name = "ppt/presentation.xml"
        pres_root = etree.fromstring(data[pres_name])
        pres_root.set("embedTrueTypeFonts", "1")

        font_lst = etree.Element(f"{{{_NS_P}}}embeddedFontLst")
        emb = etree.SubElement(font_lst, f"{{{_NS_P}}}embeddedFont")
        font_el = etree.SubElement(emb, f"{{{_NS_P}}}font")
        font_el.set("typeface", FONT)
        reg_el = etree.SubElement(emb, f"{{{_NS_P}}}regular")
        reg_el.set(f"{{{_NS_R}}}id", rid_font_reg)
        bold_el = etree.SubElement(emb, f"{{{_NS_P}}}bold")
        bold_el.set(f"{{{_NS_R}}}id", rid_font_semi)

        # embeddedFontLst по схеме идёт после sldSz/notesSz
        anchor = None
        for tag in ("notesSz", "sldSz", "sldIdLst"):
            found = pres_root.find(f"{{{_NS_P}}}{tag}")
            if found is not None:
                anchor = found
                break
        if anchor is not None:
            anchor.addnext(font_lst)
        else:
            pres_root.append(font_lst)
        data[pres_name] = etree.tostring(pres_root, xml_declaration=True,
                                         encoding="UTF-8", standalone=True)

        # --- 3. [Content_Types].xml: Default для fntdata ---
        ct_name = "[Content_Types].xml"
        ct_root = etree.fromstring(data[ct_name])
        has_fntdata = any(
            d.get("Extension") == "fntdata"
            for d in ct_root.findall(f"{{{_NS_CT}}}Default")
        )
        if not has_fntdata:
            d = etree.SubElement(ct_root, f"{{{_NS_CT}}}Default")
            d.set("Extension", "fntdata")
            d.set("ContentType", "application/x-fontdata")
        data[ct_name] = etree.tostring(ct_root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)

        # --- 4. Данные шрифтов ---
        data["ppt/fonts/font1.fntdata"] = FONT_REGULAR.read_bytes()
        semi = FONT_SEMIBOLD if FONT_SEMIBOLD.exists() else FONT_REGULAR
        data["ppt/fonts/font2.fntdata"] = semi.read_bytes()

        # --- 5. Перезаписываем zip ---
        tmp_path = pptx_path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(tmp_path, "w", zipfile.ZIP_DEFLATED) as zout:
            for name, payload in data.items():
                zout.writestr(name, payload)
        shutil.move(str(tmp_path), str(pptx_path))
        logger.info("Calleo fonts embedded into %s", pptx_path.name)

    except Exception as e:  # noqa: BLE001 — не роняем экспорт из-за встраивания
        logger.warning("Font embedding skipped (%s): %s", pptx_path.name, e)


def _next_rid(existing: set) -> str:
    """Возвращает свободный rId, не пересекающийся с existing."""
    i = 1
    while f"rId{i}" in existing:
        i += 1
    return f"rId{i}"


# ─── Цвета ────────────────────────────────────────────────────────────────────
DARK = RGBColor(0x0F, 0x11, 0x15)
GREY = RGBColor(0x55, 0x55, 0x55)

# ─── Структура блоков и подзаголовки вопросов ─────────────────────────────────
# (номера вопросов совпадают с alembic 0003_update_questions)
BLOCKS: list[tuple[str, list[int]]] = [
    ("Блок 1. Характер и первое впечатление", [1, 2]),
    ("Блок 2. Hard Skills — обучение и компетенции", [3, 4, 5]),
    ("Блок 3. Soft Skills — команда и кризисы", [6, 7, 8]),
    ("Блок 4. Клубная жизнь и тусовка", [9, 10, 11]),
    ("Блок 5. Сюжет и реальность", [12, 13]),
]

# Полный текст вопросов (совпадает с жирными заголовками из alembic
# 0003_update_questions). Именно он выводится в отчёте — жирным, а под ним
# обычным шрифтом идёт ответ педагога.
QUESTION_LABELS: dict[int, str] = {
    1: "Какой у ребёнка характер и как он обычно проявлялся в группе?",
    2: "Что самое быстрое и яркое приходит в голову про этого ребёнка?",
    3: "Стартовый уровень vs Итог: с чем ребёнок пришёл на мастер-классы и какой рывок совершил?",
    4: "Момент инсайта (озарение)",
    5: "Применение знаний",
    6: "Тип коммуникации",
    7: "Поведение в кризисе (провал миссии)",
    8: "Лидерство и ответственность",
    9: "Клубная стратегия",
    10: "Внутренняя валюта (трата зарплаты)",
    11: "Энергия",
    12: "Вера в легенду",
    13: "Неожиданное качество",
}


def _block_title_runs(block_title: str) -> list:
    """
    Готовит run'ы для заголовка блока (всё ЧЁРНЫМ):
      • убираем префикс «Блок N.»
      • часть до « — » — ПОЛУЖИРНЫМ (напр. «HARD SKILLS»)
      • часть после « — » — обычным (напр. «— ОБУЧЕНИЕ И КОМПЕТЕНЦИИ»)
    """
    title = re.sub(r"^\s*Блок\s+\d+[.\:]\s*", "", block_title).strip()
    if " — " in title:
        name, desc = title.split(" — ", 1)
        return [
            (name.upper() + " ", True, SIZE_BLOCK, DARK),
            ("— " + desc.upper(), False, SIZE_BLOCK, DARK),
        ]
    return [(title.upper(), True, SIZE_BLOCK, DARK)]


FINAL_MARKER = "=== ИТОГОВЫЙ ОТЧЁТ ==="



def _extract_final_report(text: str) -> str:
    """Возвращает итоговый отчёт (часть после маркера) либо весь текст."""
    if not text:
        return ""
    idx = text.find(FINAL_MARKER)
    if idx != -1:
        return text[idx + len(FINAL_MARKER):].strip()
    return text.strip()


# Run = (текст, bold, размер_pt, RGBColor|None)
Run = tuple[str, bool, int, "RGBColor | None"]
# Paragraph = (list[Run], выравнивание)
Para = tuple[list, "PP_ALIGN | None"]


def _estimate_height(paragraphs: list, width_emu: int) -> int:
    """Приблизительная высота набора параграфов (EMU). Намеренно с запасом."""
    width_pt = max(1.0, width_emu / 12700)
    total_pt = 2.0  # внутренние поля фрейма (сведены к минимуму)
    for runs, _align in paragraphs:
        text = "".join(r[0] for r in runs) or ""
        size = max((r[2] for r in runs), default=SIZE_BODY)
        # 0.72 — реальная средняя ширина кириллического символа (11pt Calleo
        # вмещает ~55 симв/строку, а не 90). НАМЕРЕННО с запасом: лучше чуть
        # переоценить высоту, чем обрезать текст и сломать вёрстку.
        char_pt = size * 0.72
        cpl = max(1, int(width_pt / char_pt))
        segments = text.split("\n") if text else [""]
        lines = 0
        for seg in segments:
            lines += max(1, -(-len(seg) // cpl))  # ceil
        # 1.42 — межстрочный интервал с запасом; 3.0 — отступ между параграфами.
        total_pt += lines * size * 1.42 + 3.0
    return int(total_pt * 12700) + Inches(0.06)



# Расширенный параграф с явными интервалами (в pt):
# (runs, align, space_before_pt, space_after_pt)
ParaExt = tuple[list, "PP_ALIGN | None", float, float]


def _estimate_block_height(paras_ext: list, width_emu: int) -> int:
    """
    Высота цельного блока (заголовок + пары «вопрос/ответ») в EMU с учётом
    ЯВНЫХ интервалов space_before/space_after каждого параграфа. Так весь блок
    рисуется в ОДНОМ текст-фрейме, а межстрочные и межвопросные отступы задаются
    свойствами параграфов PowerPoint → расстояния между вопросами одинаковые
    независимо от длины текста (устраняет «где-то есть отступ, где-то нет»).
    """
    width_pt = max(1.0, width_emu / 12700)
    total_pt = 6.0  # поля фрейма (с запасом)
    for runs, _align, sb_pt, sa_pt in paras_ext:
        text = "".join(r[0] for r in runs) or ""
        size = max((r[2] for r in runs), default=SIZE_BODY)
        # 0.72 — реальная ширина кириллицы (см. _estimate_height). Намеренно с
        # запасом, чтобы текст НЕ вылезал за пределы фрейма.
        char_pt = size * 0.72
        cpl = max(1, int(width_pt / char_pt))
        lines = 0
        for seg in (text.split("\n") if text else [""]):
            lines += max(1, -(-len(seg) // cpl))  # ceil
        total_pt += sb_pt + lines * size * 1.42 + sa_pt
    # запас снизу, чтобы разделитель ниже не наезжал на последнюю строку
    return int(total_pt * 12700) + Inches(0.08)



# Интервалы для блока вопросов (pt): между вопросами и после заголовка.
GAP_QUESTION_PT = 10.0   # отступ ПЕРЕД каждым вопросом (одинаковый для всех)
GAP_ANSWER_PT = 1.0      # отступ между вопросом и его ответом
GAP_HEADER_PT = 4.0      # отступ после заголовка блока


class _Deck:

    """Управляет слайдами и «курсором» вертикальной укладки контента."""

    def __init__(self, prs: Presentation) -> None:
        self.prs = prs
        self.W = prs.slide_width
        self.H = prs.slide_height
        self.m_left = Inches(0.95)
        self.m_right = Inches(0.95)
        self.m_top = Inches(0.5)
        self.m_bottom = Inches(0.5)

        self.content_w = self.W - self.m_left - self.m_right
        self.layout = self._blank_layout()
        self.slide = None
        self.y = self.m_top

    # ── служебное ──────────────────────────────────────────────────────────
    def _blank_layout(self):
        for layout in self.prs.slide_layouts:
            if len(layout.placeholders) == 0:
                return layout
        try:
            return self.prs.slide_layouts[6]
        except IndexError:
            return self.prs.slide_layouts[-1]

    @property
    def bottom_limit(self) -> int:
        return self.H - self.m_bottom

    def new_page(self, with_arrow: bool = True) -> None:
        self.slide = self.prs.slides.add_slide(self.layout)
        # чистим случайные плейсхолдеры
        for ph in list(self.slide.placeholders):
            ph._element.getparent().remove(ph._element)
        self.y = self.m_top
        if with_arrow and IMG_ARROW.exists():
            self._add_arrow()

    def _add_arrow(self) -> None:
        # Треугольник-стрелка уменьшен в 3 раза (был 0.7").
        w = Inches(0.7 / 3)
        pic = self.slide.shapes.add_picture(str(IMG_ARROW), 0, 0, width=w)
        pic.left = int(self.W - pic.width - Inches(0.25))
        pic.top = int(self.H - pic.height - Inches(0.25))


    def ensure(self, needed: int) -> None:
        """Если контент не влезает — новая страница (с стрелкой)."""
        if self.y + needed > self.bottom_limit:
            self.new_page(with_arrow=True)

    # ── картинки ──────────────────────────────────────────────────────────
    def add_image_fullwidth(self, path: Path, gap_before: int = 0,
                            gap_after: int = 0) -> None:
        if not path.exists():
            logger.warning("Asset missing: %s", path)
            return
        self.y += gap_before
        pic = self.slide.shapes.add_picture(str(path), 0, int(self.y), width=self.W)
        self.y = pic.top + pic.height + gap_after

    def add_image_centered(self, path: Path, width: int, gap_before: int = 0,
                           gap_after: int = 0) -> None:
        if not path.exists():
            logger.warning("Asset missing: %s", path)
            return
        self.y += gap_before
        left = int((self.W - width) / 2)
        pic = self.slide.shapes.add_picture(str(path), left, int(self.y), width=width)
        self.y = pic.top + pic.height + gap_after

    def add_logo_top(self, path: Path, gap_after: int = 0) -> None:
        """Лого СТРОГО в начале страницы, растянуто на всю ширину страницы
        (с сохранением пропорций — задаём только width)."""
        if not path.exists():
            logger.warning("Asset missing: %s", path)
            return
        pic = self.slide.shapes.add_picture(str(path), 0, 0, width=self.W)
        pic.left = 0
        pic.top = 0
        self.y = pic.height + gap_after

    def add_logo_bottom(self, path: Path) -> None:
        """Лого СТРОГО у нижнего края страницы, растянуто на всю ширину
        (с сохранением пропорций). Прижимается к самому низу."""
        if not path.exists():
            logger.warning("Asset missing: %s", path)
            return
        pic = self.slide.shapes.add_picture(str(path), 0, 0, width=self.W)
        pic.left = 0
        pic.top = int(self.H - pic.height)
        self.y = self.H


    # ── текст ─────────────────────────────────────────────────────────────
    def add_paragraphs(self, paragraphs: list, left: int = None, width: int = None,
                       gap_before: int = 0, gap_after: int = 0,
                       height: int = None) -> None:
        if left is None:
            left = self.m_left
        if width is None:
            width = self.content_w
        self.y += gap_before
        est = height if height is not None else _estimate_height(paragraphs, width)
        box = self.slide.shapes.add_textbox(int(left), int(self.y), int(width), int(est))
        tf = box.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        for i, (runs, align) in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align is not None:
                p.alignment = align
            for text, bold, size, rgb in runs:
                r = p.add_run()
                r.text = text
                r.font.name = FONT
                r.font.bold = bold
                r.font.size = Pt(size)
                r.font.color.rgb = rgb if rgb is not None else DARK
        self.y += est + gap_after

    def add_flow_block(self, paras_ext: list, gap_before: int = 0,
                       gap_after: int = 0) -> int:
        """
        Рисует ЦЕЛЬНЫЙ блок (заголовок + вопросы/ответы) одним текст-фреймом.
        Интервалы между параграфами задаются через space_before/space_after (pt),
        поэтому промежутки между вопросами получаются ОДИНАКОВЫМИ (PowerPoint сам
        раскладывает параграфы, а не мы позиционируем каждый бокс по «на глаз»
        оценённой высоте). Возвращает использованную высоту (EMU).
        """
        self.y += gap_before
        est = _estimate_block_height(paras_ext, self.content_w)
        box = self.slide.shapes.add_textbox(int(self.m_left), int(self.y),
                                             int(self.content_w), int(est))
        tf = box.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        for i, (runs, align, sb_pt, sa_pt) in enumerate(paras_ext):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align is not None:
                p.alignment = align
            p.space_before = Pt(sb_pt)
            p.space_after = Pt(sa_pt)
            for text, bold, size, rgb in runs:
                r = p.add_run()
                r.text = text
                r.font.name = FONT
                r.font.bold = bold
                r.font.size = Pt(size)
                r.font.color.rgb = rgb if rgb is not None else DARK
        self.y += est + gap_after
        return est


class PptxService:

    def __init__(self) -> None:
        self.reports_dir = Path(settings.reports_dir)
        self.reports_dir.mkdir(parents=True, exist_ok=True)

    def _new_presentation(self) -> Presentation:
        """
        Чистая презентация A4-портрет (вертикальный документ), собираемая с нуля.

        Намеренно НЕ наследуемся от report_template.pptx: если удалить только
        slide-id из sldIdLst, сами slide-part'ы остаются в пакете и при
        добавлении новых слайдов возникают дубли имён (ppt/slides/slide1.xml),
        что портит файл в PowerPoint. Пустой Presentation гарантирует
        корректную упаковку.
        """
        prs = Presentation()
        prs.slide_width = Inches(8.27)    # A4 портрет
        prs.slide_height = Inches(11.69)
        return prs



    def generate(
        self,
        report: Report,
        student: Student,
        shift: Shift,
        teacher: User,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        prs = self._new_presentation()
        deck = _Deck(prs)

        # Департамент берётся из РЕБЁНКА (student.department_number проставляется
        # в хендлере), т.к. shift.department_id теперь nullable/0 — смена
        # охватывает все департаменты. Фолбэк — на смену.
        dep_number = getattr(student, "department_number", None)
        if dep_number is None:
            dep_number = shift.department_id or 0
        dep_rgb = _hex_to_rgb(get_department_color(dep_number))
        dep_name = get_department_name(dep_number)


        answers = q_answers or parse_numbered_answers(report.generated_text or "")
        final_report = _extract_final_report(report.generated_text or "")

        # ── СТРАНИЦА 1 ────────────────────────────────────────────────────
        deck.new_page(with_arrow=True)
        # Заглавное лого — СТРОГО в начале страницы, на всю ширину (с сохранением
        # пропорций).
        deck.add_logo_top(IMG_LOGO_TOP, gap_after=Inches(0.15))
        # Верхний разделитель тоже «обволакивает» профиль: профиль сверху ЗАХОДИТ
        # внутрь separator_corner (отрицательный gap_after).
        overlap = Inches(0.14)
        deck.add_image_fullwidth(IMG_SEP_CORNER, gap_after=-overlap)

        self._add_profile(

            deck,
            shift_name=shift.name or "",
            shift_dates=_resolve_shift_dates(shift),
            dep_name=dep_name,
            dep_rgb=dep_rgb,
            student_name=student.full_name or "",
        )

        # Нижний разделитель «обнимает» профиль И легенду: профиль сверху ЗАХОДИТ
        # в линию (отрицательный gap_before), а заголовок ЛЕГЕНДЫ снизу — тоже
        # (отрицательный gap_after). Так ФИО/департамент «влипают» в разделитель.
        deck.add_image_fullwidth(IMG_SEP_LINE, gap_before=-overlap,

                                 gap_after=-overlap)

        # Легенда смены — заголовок ЧЁРНЫЙ (не в цвет департамента).
        deck.add_paragraphs(
            [([("ЛЕГЕНДА СМЕНЫ", True, SIZE_BLOCK, DARK)], PP_ALIGN.LEFT)],
            gap_after=Inches(0.02),
        )

        deck.add_paragraphs(
            [([((shift_context or "Контекст смены не указан."), False, SIZE_BODY, DARK)],
              PP_ALIGN.LEFT)],
        )

        # ── СТРАНИЦЫ БЛОКОВ ───────────────────────────────────────────────
        # Разделитель ОТКРЫВАЕТ и ЗАКРЫВАЕТ каждый блок и «ОБВОЛАКИВАЕТ» текст:
        # контент чуть ЗАХОДИТ внутрь картинки разделителя (отрицательный
        # перехлёст OVERLAP), но не по центру — только слегка. Так линия сверху
        # и снизу как бы обнимают блок.

        deck.new_page(with_arrow=True)
        deck.add_image_fullwidth(IMG_SEP_LINE, gap_after=-overlap)
        last_idx = len(BLOCKS) - 1
        for b_idx, (block_title, q_nums) in enumerate(BLOCKS):
            self._add_block(deck, block_title, q_nums, answers, dep_rgb)
            # После ПОСЛЕДНЕГО блока разделитель НЕ ставим.
            if b_idx != last_idx:
                deck.add_image_fullwidth(IMG_SEP_LINE, gap_before=-overlap,
                                         gap_after=-overlap)




        # ── ПОСЛЕДНЯЯ СТРАНИЦА (без стрелки) ──────────────────────────────
        deck.new_page(with_arrow=False)
        # Заголовок «ОБЩЕЕ ЗАКЛЮЧЕНИЕ» так же «обволакивается» разделителями:
        # чуть заходит внутрь верхнего (first_last) и нижнего (last) разделителя.
        deck.add_image_fullwidth(IMG_SEP_FIRST_LAST, gap_after=-overlap)
        # Заголовок «ОБЩЕЕ ЗАКЛЮЧЕНИЕ» — ЧЁРНЫЙ (не в цвет департамента).
        deck.add_paragraphs(
            [([("ОБЩЕЕ ЗАКЛЮЧЕНИЕ ПЕДАГОГА", True, SIZE_BLOCK, DARK)], PP_ALIGN.LEFT)],
            gap_after=Inches(0.02),
        )
        deck.add_paragraphs(
            [([((final_report or "—"), False, SIZE_BODY, DARK)], PP_ALIGN.LEFT)],
            gap_after=-overlap,
        )
        deck.add_image_fullwidth(IMG_SEP_LAST, gap_before=0,
                                 gap_after=Inches(0.2))

        # Финальное лого — СТРОГО у нижнего края страницы, на всю ширину
        # (с сохранением пропорций).
        deck.add_logo_bottom(IMG_LOGO_BOTTOM)


        output_path = self.reports_dir / safe_pptx_filename(student.full_name)
        prs.save(str(output_path))
        _embed_calleo_fonts(output_path)
        logger.info("PPTX built: %s for student=%r", output_path, student.full_name)
        return output_path

    # ── профиль (метки СЛЕВА к разделителю, значения СПРАВА от него) ────────
    def _add_profile(self, deck: _Deck, shift_name: str, shift_dates: str,
                     dep_name: str, dep_rgb: RGBColor, student_name: str) -> None:
        # Метки — приглушённо-серые, обычным; значения — ПОЛУЖИРНЫЕ и КАПСОМ.
        labels = ["Название смены", "Дата смены", "Департамент", "ФИО"]
        values = [
            ((shift_name or "—").upper(), DARK),
            ((shift_dates or "—").upper(), DARK),
            ((dep_name or "—").upper(), dep_rgb),
            ((student_name or "—").upper(), DARK),
        ]

        row_h = Inches(0.36)
        block_h = row_h * len(labels)
        top = deck.y

        # Вертикальный разделитель — чуть левее центра страницы.
        sep_x = Inches(2.75)

        # Левая колонка меток: право-выравнивание, край — чуть левее разделителя.
        label_left = deck.m_left
        label_w = sep_x - Inches(0.22) - label_left
        label_paras = [
            ([(lbl, False, SIZE_PROFILE, GREY)], PP_ALIGN.RIGHT) for lbl in labels
        ]
        lbox = deck.slide.shapes.add_textbox(int(label_left), int(top),
                                             int(label_w), int(block_h))
        self._fill_rows(lbox, label_paras, row_h)

        # Значения: ВЫРОВНЕНЫ ПО РАЗДЕЛИТЕЛЮ (слева-направо, сразу после него),
        # полужирным. Так они «висят» на разделителе, как в макете.
        value_left = sep_x + Inches(0.22)
        value_w = deck.W - deck.m_right - value_left
        value_paras = [
            ([(val, True, SIZE_PROFILE, rgb)], PP_ALIGN.LEFT) for val, rgb in values
        ]
        vbox = deck.slide.shapes.add_textbox(int(value_left), int(top),
                                             int(value_w), int(block_h))
        self._fill_rows(vbox, value_paras, row_h)

        # Вертикальный разделитель.
        if IMG_SEP_PROFILE_V.exists():
            pic = deck.slide.shapes.add_picture(
                str(IMG_SEP_PROFILE_V), 0, 0, height=int(block_h))
            pic.left = int(sep_x - pic.width / 2)
            pic.top = int(top)

        deck.y = top + block_h



    @staticmethod
    def _fill_rows(box, paragraphs: list, row_h: int) -> None:
        tf = box.text_frame
        tf.word_wrap = True
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for i, (runs, align) in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if align is not None:
                p.alignment = align
            # равномерно распределяем строки по высоте row_h
            p.space_after = Pt(max(0, (row_h / 12700) - SIZE_PROFILE * 1.2))
            for text, bold, size, rgb in runs:
                r = p.add_run()
                r.text = text
                r.font.name = FONT
                r.font.bold = bold
                r.font.size = Pt(size)
                r.font.color.rgb = rgb if rgb is not None else DARK

    # ── блок с вопросами ──────────────────────────────────────────────────
    def _add_block(self, deck: _Deck, block_title: str, q_nums: list[int],
                   answers: dict, dep_rgb: RGBColor) -> None:
        """
        Весь блок (заголовок + все пары «вопрос/ответ») рисуется ОДНИМ текст-
        фреймом через deck.add_flow_block. Отступы между вопросами задаются
        одинаковым space_before (GAP_QUESTION_PT), поэтому промежутки ровные.
        """
        paras: list = []
        # Заголовок блока — ЧЁРНЫЙ: имя блока полужирным, описание обычным.
        paras.append((_block_title_runs(block_title), PP_ALIGN.LEFT,
                      0.0, GAP_HEADER_PT))
        for qn in q_nums:
            label = QUESTION_LABELS.get(qn, f"Вопрос {qn}")
            answer = answers.get(qn) or answers.get(str(qn)) or "—"
            # Вопрос — ПОЛУЖИРНЫМ, единый отступ сверху (равные промежутки).
            paras.append(([(f"{qn}. {label}", True, SIZE_QUESTION, DARK)],
                          PP_ALIGN.LEFT, GAP_QUESTION_PT, GAP_ANSWER_PT))
            # Ответ педагога — обычным, вплотную к своему вопросу.
            paras.append(([(str(answer), False, SIZE_BODY, DARK)],
                          PP_ALIGN.LEFT, 0.0, 0.0))

        est = _estimate_block_height(paras, deck.content_w)
        deck.ensure(est)
        deck.add_flow_block(paras)


    # ── конвертация в PDF (через LibreOffice/soffice) ──────────────────────
    def _to_pdf(self, pptx_path: Path) -> Path:
        """
        Конвертирует .pptx → .pdf через headless LibreOffice.
        Требует установленного soffice (см. Dockerfile). Встроенные шрифты
        Calleo подхватываются автоматически.
        """
        outdir = pptx_path.parent
        subprocess.run(
            [
                "soffice", "--headless", "--nologo", "--nofirststartwizard",
                "--convert-to", "pdf", "--outdir", str(outdir), str(pptx_path),
            ],
            check=True,
            timeout=180,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
        )
        pdf_path = pptx_path.with_suffix(".pdf")
        if not pdf_path.exists():
            raise RuntimeError(f"PDF не создан: {pdf_path}")
        logger.info("PDF built: %s", pdf_path.name)
        return pdf_path

    def generate_pdf(
        self,
        report: Report,
        student: Student,
        shift: Shift,
        teacher: User,
        q_answers: dict[int, str] | None = None,
        shift_context: str | None = None,
    ) -> Path:
        """Собирает PPTX и конвертирует его в PDF, возвращает путь к PDF."""
        pptx_path = self.generate(
            report=report, student=student, shift=shift, teacher=teacher,
            q_answers=q_answers, shift_context=shift_context,
        )
        return self._to_pdf(pptx_path)




